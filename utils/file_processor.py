"""
File processing utilities for AI Approach Chat.
Handles file upload, validation, text extraction, and chunking.
"""
import os
import json
import uuid
import fitz  # PyMuPDF
import docx
import pptx
import pdfplumber
from datetime import datetime
from pathlib import Path

class FileProcessor:
    def __init__(self, base_path):
        """Initialize the file processor with base paths for storage."""
        self.base_path = base_path
        self.files_dir = os.path.join(base_path, 'data', 'files')
        self.chunks_dir = os.path.join(base_path, 'data', 'chunks')
        self.metadata_path = os.path.join(base_path, 'data', 'files', 'metadata.json')
        
        # Ensure directories exist
        os.makedirs(self.files_dir, exist_ok=True)
        os.makedirs(self.chunks_dir, exist_ok=True)
        
        # Initialize or load metadata
        if os.path.exists(self.metadata_path):
            with open(self.metadata_path, 'r') as f:
                self.metadata = json.load(f)
        else:
            self.metadata = {}
            self._save_metadata()

    def _save_metadata(self):
        """Save metadata to JSON file."""
        with open(self.metadata_path, 'w') as f:
            json.dump(self.metadata, f, indent=2)

    def _generate_file_id(self):
        """Generate a unique file ID."""
        return str(uuid.uuid4())

    def _get_file_extension(self, filename):
        """Get the file extension."""
        return os.path.splitext(filename)[1].lower()

    def validate_file(self, file_path):
        """
        Validate if the file contains extractable text.
        Returns (is_valid, message)
        """
        ext = self._get_file_extension(file_path)
        
        try:
            if ext == '.pdf':
                # Check if PDF has extractable text
                with fitz.open(file_path) as doc:
                    text_content = ""
                    for page in doc:
                        text_content += page.get_text()
                    
                    # Also try with pdfplumber for scanned PDFs
                    if not text_content.strip():
                        with pdfplumber.open(file_path) as pdf:
                            for page in pdf.pages:
                                text_content += page.extract_text() or ""
                    
                    if not text_content.strip():
                        return False, "The PDF file contains no extractable text. It might be a scanned document."
            
            elif ext == '.docx':
                # Check if DOCX has content
                doc = docx.Document(file_path)
                if len(doc.paragraphs) == 0:
                    return False, "The Word document appears to be empty."
            
            elif ext == '.pptx':
                # Check if PPTX has content
                prs = pptx.Presentation(file_path)
                if len(prs.slides) == 0:
                    return False, "The PowerPoint presentation appears to be empty."
            
            else:
                return False, f"Unsupported file format: {ext}"
            
            return True, "File is valid and contains extractable text."
        
        except Exception as e:
            return False, f"Error validating file: {str(e)}"

    def process_file(self, file_path, original_filename, username):
        """
        Process a file: validate, save, extract text, chunk content.
        Returns (success, file_id or error_message)
        """
        # Validate file
        is_valid, message = self.validate_file(file_path)
        if not is_valid:
            return False, message
        
        # Generate file ID and save file
        file_id = self._generate_file_id()
        ext = self._get_file_extension(original_filename)
        saved_path = os.path.join(self.files_dir, f"{file_id}{ext}")
        
        # Copy file to storage location
        with open(file_path, 'rb') as src, open(saved_path, 'wb') as dst:
            dst.write(src.read())
        
        # Extract and chunk text
        chunks = self._extract_and_chunk_text(saved_path, file_id, ext)
        
        if not chunks:
            # If no chunks were created (extraction failed)
            os.remove(saved_path)
            return False, "Failed to extract text from the file."
        
        # Save chunks
        chunks_path = os.path.join(self.chunks_dir, f"{file_id}.json")
        with open(chunks_path, 'w') as f:
            json.dump(chunks, f, indent=2)
        
        # Update metadata
        file_info = {
            'file_id': file_id,
            'filename': original_filename,
            'extension': ext,
            'upload_date': datetime.now().isoformat(),
            'username': username,
            'file_path': saved_path,
            'chunks_path': chunks_path,
            'chunk_count': len(chunks)
        }
        
        if username not in self.metadata:
            self.metadata[username] = []
        
        self.metadata[username].append(file_info)
        self._save_metadata()
        
        return True, file_id

    def _extract_text_from_pdf(self, file_path):
        """Extract text from PDF file with page numbers."""
        text_by_page = []
        
        # Try PyMuPDF first
        try:
            with fitz.open(file_path) as doc:
                for page_num, page in enumerate(doc):
                    text = page.get_text()
                    if text.strip():
                        text_by_page.append({
                            'page': page_num + 1,
                            'text': text
                        })
        except Exception:
            # Fall back to pdfplumber
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page_num, page in enumerate(pdf.pages):
                        text = page.extract_text()
                        if text and text.strip():
                            text_by_page.append({
                                'page': page_num + 1,
                                'text': text
                            })
            except Exception as e:
                print(f"Error extracting text from PDF: {str(e)}")
        
        return text_by_page

    def _extract_text_from_docx(self, file_path):
        """Extract text from DOCX file."""
        doc = docx.Document(file_path)
        text = []
        
        # Approximate page breaks (rough estimate)
        paragraphs_per_page = 15
        current_page = 1
        current_page_text = ""
        paragraph_count = 0
        
        for para in doc.paragraphs:
            if para.text.strip():
                current_page_text += para.text + "\n"
                paragraph_count += 1
                
                if paragraph_count >= paragraphs_per_page:
                    text.append({
                        'page': current_page,
                        'text': current_page_text
                    })
                    current_page += 1
                    current_page_text = ""
                    paragraph_count = 0
        
        # Add the last page if it has content
        if current_page_text.strip():
            text.append({
                'page': current_page,
                'text': current_page_text
            })
        
        return text

    def _extract_text_from_pptx(self, file_path):
        """Extract text from PPTX file with slide numbers."""
        prs = pptx.Presentation(file_path)
        text_by_slide = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            
            if slide_text.strip():
                text_by_slide.append({
                    'page': slide_num + 1,  # Using 'page' for consistency
                    'text': slide_text
                })
        
        return text_by_slide

    def _extract_text(self, file_path, ext):
        """Extract text based on file type."""
        if ext == '.pdf':
            return self._extract_text_from_pdf(file_path)
        elif ext == '.docx':
            return self._extract_text_from_docx(file_path)
        elif ext == '.pptx':
            return self._extract_text_from_pptx(file_path)
        return []

    def _chunk_text(self, text_by_page, file_id):
        """
        Chunk text into meaningful segments.
        Smart chunking to maintain context and coherence.
        """
        chunks = []
        chunk_id = 0
        
        for page_data in text_by_page:
            page_num = page_data['page']
            text = page_data['text']
            
            # Split text into paragraphs
            paragraphs = text.split('\n\n')
            
            current_chunk = ""
            current_chunk_size = 0
            target_chunk_size = 1000  # Target characters per chunk
            
            for paragraph in paragraphs:
                paragraph = paragraph.strip()
                if not paragraph:
                    continue
                
                # If adding this paragraph would exceed target size and we already have content,
                # save the current chunk and start a new one
                if current_chunk_size + len(paragraph) > target_chunk_size and current_chunk:
                    chunk_id += 1
                    chunks.append({
                        'chunk_id': f"{file_id}_{chunk_id}",
                        'file_id': file_id,
                        'page': page_num,
                        'content': current_chunk.strip()
                    })
                    current_chunk = paragraph + "\n"
                    current_chunk_size = len(paragraph)
                else:
                    current_chunk += paragraph + "\n"
                    current_chunk_size += len(paragraph)
            
            # Add the last chunk if it has content
            if current_chunk.strip():
                chunk_id += 1
                chunks.append({
                    'chunk_id': f"{file_id}_{chunk_id}",
                    'file_id': file_id,
                    'page': page_num,
                    'content': current_chunk.strip()
                })
        
        return chunks

    def _extract_and_chunk_text(self, file_path, file_id, ext):
        """Extract text from file and chunk it."""
        text_by_page = self._extract_text(file_path, ext)
        if not text_by_page:
            return []
        
        chunks = self._chunk_text(text_by_page, file_id)
        return chunks

    def get_file_metadata(self, file_id):
        """Get metadata for a specific file."""
        for username, files in self.metadata.items():
            for file_info in files:
                if file_info['file_id'] == file_id:
                    return file_info
        return None

    def get_user_files(self, username):
        """Get all files for a specific user."""
        if username in self.metadata:
            return self.metadata[username]
        return []

    def get_chunks(self, file_id):
        """Get all chunks for a specific file."""
        chunks_path = os.path.join(self.chunks_dir, f"{file_id}.json")
        if os.path.exists(chunks_path):
            with open(chunks_path, 'r') as f:
                return json.load(f)
        return []

    def delete_file(self, file_id, username):
        """Delete a file and its chunks."""
        if username not in self.metadata:
            return False
        
        file_info = None
        for i, info in enumerate(self.metadata[username]):
            if info['file_id'] == file_id:
                file_info = info
                self.metadata[username].pop(i)
                break
        
        if not file_info:
            return False
        
        # Delete file
        if os.path.exists(file_info['file_path']):
            os.remove(file_info['file_path'])
        
        # Delete chunks
        if os.path.exists(file_info['chunks_path']):
            os.remove(file_info['chunks_path'])
        
        self._save_metadata()
        return True
